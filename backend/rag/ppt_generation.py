import hashlib
import os
import re
import time
import uuid
from typing import Any, Dict, List, Optional, Sequence, Tuple

import qrcode
import requests
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from werkzeug.utils import secure_filename

from backend.extensions import db
from backend.models.material import Material


PPT_EXPORT_FILENAME_PREFIX = "lesson-ppt"
PPT_THEMES = {"clean", "tech", "vivid"}
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def normalize_ppt_outline_to_12_slides(spec: Dict[str, Any], target_slide_count: int = 12) -> Dict[str, Any]:
    requirement = spec.get("requirement_summary") if isinstance(spec.get("requirement_summary"), dict) else {}
    ppt_outline = spec.get("ppt_outline") if isinstance(spec.get("ppt_outline"), list) else []
    docx_outline = spec.get("docx_outline") if isinstance(spec.get("docx_outline"), list) else []

    content_slides: List[Dict[str, Any]] = []
    summary_candidates: List[Dict[str, Any]] = []

    for item in ppt_outline:
        if not isinstance(item, dict):
            continue
        slide_type = str(item.get("slide_type") or "").strip().lower()
        if slide_type in {"summary", "conclusion", "ending", "tail", "end"}:
            summary_candidates.append(item)
            continue
        if slide_type in {"cover", "title", "toc", "agenda", "catalog"}:
            continue
        content_slides.append(_normalize_outline_slide(item))

    content_slides = [item for item in content_slides if item["title"] or item["bullets"]]
    content_slides = _split_dense_content_slides(content_slides, desired_count=min(8, max(1, target_slide_count - 4)))

    if len(content_slides) < 8:
        for item in _build_supplemental_content_slides(requirement, docx_outline):
            if len(content_slides) >= 8:
                break
            duplicate = any(existing["title"] == item["title"] and existing["bullets"] == item["bullets"] for existing in content_slides)
            if not duplicate:
                content_slides.append(item)

    while len(content_slides) < 8:
        index = len(content_slides) + 1
        content_slides.append(
            {
                "slide_type": "content",
                "title": f"Core Topic {index}",
                "goal": "",
                "bullets": ["Explain the core concept clearly", "Connect the concept to classroom practice"],
                "visual_suggestion": "",
                "source_refs": [],
            }
        )

    content_slides = content_slides[:8]
    topic = str(requirement.get("topic") or requirement.get("chapter_title") or requirement.get("grade_subject") or "AI Lesson Deck").strip()
    subtitle = " | ".join(
        [
            item
            for item in [
                str(requirement.get("grade_subject") or "").strip(),
                str(requirement.get("chapter_title") or "").strip(),
                str(requirement.get("duration") or "").strip(),
            ]
            if item
        ]
    )

    return {
        "cover": {"title": topic or "AI Lesson Deck", "subtitle": subtitle},
        "toc": {"items": [item.get("toc_title") or item["title"] or f"Topic {index + 1}" for index, item in enumerate(content_slides)]},
        "content_slides": content_slides,
        "summary": _build_summary_slide(summary_candidates, requirement, content_slides),
        "ending": {"title": "Thank You", "subtitle": "易度新星 EduNova generated presentation"},
    }


PPT_TEMPLATE_ROOT = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "templates", "lesson_plan", "ppt",
)
PPT_TEMPLATE_MANIFEST = os.path.join(PPT_TEMPLATE_ROOT, "manifest.json")


def resolve_ppt_template(template_profile: Optional[str] = None) -> Dict[str, Any]:
    profile = str(template_profile or "default").strip().lower() or "default"
    manifest = _load_ppt_manifest()
    profiles = manifest.get("profiles") if isinstance(manifest.get("profiles"), dict) else {}
    profile_meta = profiles.get(profile) if isinstance(profiles.get(profile), dict) else {}
    if not profile_meta and profile != "default":
        profile_meta = profiles.get("default") if isinstance(profiles.get("default"), dict) else {}
        profile = "default"

    template_name = str(profile_meta.get("template") or "official.pptx").strip()
    template_path = os.path.join(PPT_TEMPLATE_ROOT, template_name)
    if not os.path.exists(template_path):
        template_path = None

    layout_map = profile_meta.get("layout_map") if isinstance(profile_meta.get("layout_map"), dict) else {}

    return {
        "profile": profile,
        "template_path": template_path,
        "target_slide_count": int(profile_meta.get("target_slide_count", 12) or 12),
        "layout_map": layout_map,
    }


def _load_ppt_manifest() -> Dict[str, Any]:
    if not os.path.exists(PPT_TEMPLATE_MANIFEST):
        return {}
    try:
        import json
        with open(PPT_TEMPLATE_MANIFEST, "r", encoding="utf-8") as f:
            data = json.load(f)
        return data if isinstance(data, dict) else {}
    except Exception:
        return {}


def build_game_entry(
    upload_root: str,
    course_id: int,
    title: str,
    material_title: str,
    download_url: str,
) -> Dict[str, Any]:
    qr_assets_dir = os.path.join(upload_root, "materials", str(course_id), ".generated_assets")
    os.makedirs(qr_assets_dir, exist_ok=True)
    qr_path = _build_qr_code_image(download_url, qr_assets_dir)
    return {
        "title": str(title or "小游戏入口页").strip() or "小游戏入口页",
        "subtitle": "扫码或点击链接下载并运行课堂互动小游戏",
        "material_title": str(material_title or "lesson-game.html").strip() or "lesson-game.html",
        "download_url": str(download_url or "").strip(),
        "teacher_tip": "建议先将 HTML 文件下载到本机，再双击打开进行课堂演示或课后练习。",
        "qr_image_path": qr_path,
    }


def build_theme_config(spec: Dict[str, Any], requested_theme: Optional[str] = None) -> Dict[str, Any]:
    theme = str(requested_theme or "").strip().lower()
    if theme not in PPT_THEMES:
        theme = _infer_theme_from_spec(spec)

    if theme == "tech":
        return {
            "name": "tech",
            "font_family": "Microsoft YaHei",
            "bg": (11, 18, 33),
            "text": (237, 245, 255),
            "muted": (155, 185, 220),
            "accent": (35, 214, 220),
            "accent_soft": (34, 69, 122),
            "cover_band": (28, 46, 80),
            "summary_bg": (18, 31, 56),
            "footer": (120, 157, 214),
            "title_size": 28,
            "cover_title_size": 30,
            "body_size": 17,
            "small_size": 10,
        }
    if theme == "vivid":
        return {
            "name": "vivid",
            "font_family": "Microsoft YaHei",
            "bg": (255, 248, 242),
            "text": (72, 39, 34),
            "muted": (138, 99, 84),
            "accent": (242, 104, 56),
            "accent_soft": (255, 217, 193),
            "cover_band": (255, 236, 222),
            "summary_bg": (255, 243, 233),
            "footer": (186, 93, 59),
            "title_size": 28,
            "cover_title_size": 30,
            "body_size": 17,
            "small_size": 10,
        }
    return {
        "name": "clean",
        "font_family": "Microsoft YaHei",
        "bg": (247, 249, 252),
        "text": (37, 47, 64),
        "muted": (98, 112, 137),
        "accent": (40, 99, 210),
        "accent_soft": (220, 232, 252),
        "cover_band": (231, 238, 250),
        "summary_bg": (240, 245, 252),
        "footer": (107, 125, 154),
        "title_size": 28,
        "cover_title_size": 30,
        "body_size": 17,
        "small_size": 10,
    }


def collect_image_candidates(
    slide: Dict[str, Any],
    processed_sources: Sequence[Dict[str, Any]],
    upload_root: str,
    used_candidates: Optional[set] = None,
) -> List[Dict[str, Any]]:
    used_candidates = used_candidates or set()
    keywords = _extract_slide_keywords(slide)
    candidates: List[Dict[str, Any]] = []

    for source in processed_sources:
        if not isinstance(source, dict):
            continue
        kind = str(source.get("kind") or "").strip().lower()
        mapping = source.get("mapping") if isinstance(source.get("mapping"), dict) else {}
        usage = str(mapping.get("usage") or "").strip().lower()
        knowledge_point = str(mapping.get("knowledge_point") or "").strip()
        base_text = " ".join(
            [
                str(mapping.get("file_name") or "").strip(),
                knowledge_point,
                str(source.get("summary") or "").strip(),
                str(source.get("raw_text") or "").strip()[:180],
            ]
        )

        if kind == "image":
            image_path = _resolve_upload_path(upload_root, mapping.get("file_path"))
            if image_path and image_path not in used_candidates:
                score = 60 + _score_keyword_match(base_text, keywords) + (22 if usage == "image_asset" else 0)
                candidates.append(_build_candidate("teacher_upload", image_path, score, knowledge_point or str(mapping.get("file_name") or "").strip(), str(mapping.get("file_name") or "").strip()))
            continue

        if kind == "ppt":
            assets = source.get("assets") if isinstance(source.get("assets"), dict) else {}
            for ppt_slide in assets.get("slides") or []:
                if not isinstance(ppt_slide, dict):
                    continue
                ppt_text = " ".join(
                    [
                        str(ppt_slide.get("title") or "").strip(),
                        str(ppt_slide.get("text") or "").strip(),
                        str(ppt_slide.get("notes") or "").strip(),
                        knowledge_point,
                    ]
                )
                for rel_path in ppt_slide.get("image_paths") or []:
                    image_path = _resolve_upload_path(upload_root, rel_path)
                    if not image_path or image_path in used_candidates:
                        continue
                    score = 38 + _score_keyword_match(ppt_text, keywords) + (12 if _slide_ref_matches(slide, mapping) else 0)
                    candidates.append(_build_candidate("ppt_embedded", image_path, score, str(ppt_slide.get("title") or "").strip(), str(mapping.get("file_name") or "").strip()))
            continue

        if kind == "video":
            assets = source.get("assets") if isinstance(source.get("assets"), dict) else {}
            chunks = source.get("chunks") if isinstance(source.get("chunks"), list) else []
            chunk_by_index = {
                int(item.get("index")): item
                for item in chunks
                if isinstance(item, dict) and str(item.get("index") or "").isdigit()
            }
            for index, frame in enumerate(assets.get("keyframes") or [], start=1):
                if not isinstance(frame, dict):
                    continue
                image_path = _resolve_upload_path(upload_root, frame.get("image_path"))
                if not image_path or image_path in used_candidates:
                    continue
                chunk = chunk_by_index.get(index, {})
                frame_text = " ".join(
                    [
                        str(frame.get("summary") or "").strip(),
                        str(frame.get("ocr_text") or "").strip(),
                        str(chunk.get("summary") or "").strip(),
                        knowledge_point,
                    ]
                )
                importance = float(chunk.get("importance_score", 0) or 0)
                score = 28 + int(importance * 20) + _score_keyword_match(frame_text, keywords)
                candidates.append(_build_candidate("video_keyframe", image_path, score, str(frame.get("summary") or frame.get("ocr_text") or "").strip(), str(mapping.get("file_name") or "").strip()))

    candidates.sort(key=lambda item: (-int(item["score"]), _landscape_bonus(item["path_or_url"])))
    return candidates


def search_gallery_images(
    slide: Dict[str, Any],
    spec: Dict[str, Any],
    course_id: int,
    upload_root: str,
    used_candidates: Optional[set] = None,
    api_key: Optional[str] = None,
) -> Tuple[List[Dict[str, Any]], List[str]]:
    used_candidates = used_candidates or set()
    api_key = str(api_key or os.getenv("PEXELS_API_KEY") or "").strip()
    if not api_key:
        return [], ["PEXELS_API_KEY is not configured, gallery fallback skipped"]

    queries = _build_gallery_queries(slide, spec)
    headers = {"Authorization": api_key}

    for query in queries:
        try:
            response = requests.get(
                "https://api.pexels.com/v1/search",
                headers=headers,
                params={"query": query, "per_page": 3, "orientation": "landscape"},
                timeout=10,
            )
            response.raise_for_status()
            payload = response.json()
        except Exception as exc:
            return [], [f"Gallery search failed: {str(exc)}"]

        gallery_candidates: List[Dict[str, Any]] = []
        for item in payload.get("photos") or []:
            if not isinstance(item, dict):
                continue
            photo_id = str(item.get("id") or "").strip()
            src = item.get("src") if isinstance(item.get("src"), dict) else {}
            image_url = str(src.get("large2x") or src.get("large") or src.get("original") or "").strip()
            if not photo_id or not image_url:
                continue
            image_path = _download_gallery_image(image_url, photo_id, query, course_id, upload_root)
            if not image_path or image_path in used_candidates:
                continue
            gallery_candidates.append(_build_candidate("gallery", image_path, 10 + _landscape_bonus(image_path), str(item.get("alt") or query).strip(), "pexels"))
        if gallery_candidates:
            gallery_candidates.sort(key=lambda item: (-int(item["score"]), item["path_or_url"]))
            return gallery_candidates, []

    return [], ["No gallery image matched, fallback used local assets or text-only slides"]


def render_pptx(
    normalized_plan: Dict[str, Any],
    spec: Dict[str, Any],
    theme_config: Dict[str, Any],
    slide_images: Sequence[Optional[Dict[str, Any]]],
    output_path: str,
    game_entry: Optional[Dict[str, Any]] = None,
    template_profile: Optional[str] = None,
) -> Dict[str, int]:
    template_meta = resolve_ppt_template(template_profile)
    template_path = template_meta.get("template_path")
    layout_map = template_meta.get("layout_map") if isinstance(template_meta.get("layout_map"), dict) else {}

    if template_path:
        prs = Presentation(template_path)
    else:
        prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _render_cover_slide(prs, normalized_plan.get("cover", {}), theme_config, spec, layout_map)
    _render_toc_slide(prs, normalized_plan.get("toc", {}), theme_config, layout_map)

    image_stats = {"teacher_images": 0, "keyframes": 0, "gallery_images": 0, "text_only_slides": 0}
    content_slides = normalized_plan.get("content_slides") if isinstance(normalized_plan.get("content_slides"), list) else []

    for index, slide_data in enumerate(content_slides):
        image = slide_images[index] if index < len(slide_images) else None
        layout_name = "image-left" if index % 2 == 0 else "image-right"
        if not image:
            layout_name = "text-focus"
            image_stats["text_only_slides"] += 1
        elif image["source_type"] == "teacher_upload":
            image_stats["teacher_images"] += 1
        elif image["source_type"] == "video_keyframe":
            image_stats["keyframes"] += 1
        elif image["source_type"] == "gallery":
            image_stats["gallery_images"] += 1
        _render_content_slide(prs, slide_data, theme_config, image, layout_name, index + 3, game_entry=game_entry, layout_map=layout_map)

    _render_summary_slide(prs, normalized_plan.get("summary", {}), theme_config, 11, layout_map)
    if game_entry:
        _render_game_entry_slide(prs, game_entry, theme_config, 12, layout_map)
    else:
        _render_ending_slide(prs, normalized_plan.get("ending", {}), theme_config, 12, layout_map)
    prs.save(output_path)
    return image_stats


def _get_layout(prs: Presentation, layout_map: Dict[str, Any], key: str) -> Any:
    """Return the slide layout for *key* from *layout_map*, falling back to blank."""
    idx = layout_map.get(key)
    if idx is not None:
        try:
            return prs.slide_layouts[int(idx)]
        except (IndexError, ValueError):
            pass
    return prs.slide_layouts[6]


def persist_generated_material(
    course_id: int,
    course_name: str,
    topic: str,
    theme_name: str,
    output_path: str,
    image_stats: Dict[str, int],
) -> Material:
    file_hash = _calculate_file_hash(output_path)
    if not file_hash:
        raise ValueError("Failed to calculate hash for generated PPT")

    existing = Material.query.filter_by(course_id=course_id, file_hash=file_hash).first()
    if existing and existing.file_path:
        existing_abs_path = _material_upload_to_abs_path(existing.file_path)
        if existing_abs_path and os.path.exists(existing_abs_path):
            if os.path.abspath(existing_abs_path) != os.path.abspath(output_path) and os.path.exists(output_path):
                os.remove(output_path)
            return existing

    filename = os.path.basename(output_path)
    relative_upload_path = f"/uploads/materials/{course_id}/{filename}"
    content = f"Generated PPT from lesson plan spec | course={course_name} | topic={topic} | theme={theme_name} | images={image_stats}"

    if existing:
        existing.title = filename
        existing.material_type = "PowerPoint"
        existing.file_path = relative_upload_path
        existing.content = content
        db.session.commit()
        return existing

    material = Material(
        title=filename,
        material_type="PowerPoint",
        file_path=relative_upload_path,
        file_hash=file_hash,
        content=content,
        course_id=course_id,
    )
    db.session.add(material)
    db.session.commit()
    return material


def build_output_path(upload_root: str, course_id: int, topic: str) -> str:
    materials_dir = os.path.join(upload_root, "materials", str(course_id))
    os.makedirs(materials_dir, exist_ok=True)
    base_name = secure_filename(topic) or PPT_EXPORT_FILENAME_PREFIX
    filename = f"{base_name}_{time.strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:6]}.pptx"
    return os.path.join(materials_dir, filename)


def _normalize_outline_slide(item: Dict[str, Any]) -> Dict[str, Any]:
    raw_title = str(item.get("title") or "").strip()
    bullets = _clip_bullets(item.get("bullets"), minimum=2)
    if not bullets:
        goal = str(item.get("goal") or "").strip()
        bullets = [goal] if goal else []
    if len(bullets) < 2:
        bullets.append("Use one concrete classroom example")
    title = _shorten_slide_title(raw_title, bullets, str(item.get("goal") or "").strip())
    title_detail = _extract_title_detail_as_bullet(raw_title, title)
    if title_detail:
        bullets = [title_detail] + list(bullets)
    bullets = _expand_slide_bullets(
        bullets=_prepare_display_bullets(bullets, title or raw_title),
        title=title or raw_title,
        goal=str(item.get("goal") or "").strip(),
    )
    return {
        "slide_type": str(item.get("slide_type") or "content").strip(),
        "title": title,
        "toc_title": _build_toc_title(title, raw_title),
        "goal": str(item.get("goal") or "").strip(),
        "bullets": bullets[:4],
        "visual_suggestion": str(item.get("visual_suggestion") or "").strip(),
        "source_refs": [str(value).strip() for value in item.get("source_refs", []) if str(value).strip()],
    }


def _build_supplemental_content_slides(requirement: Dict[str, Any], docx_outline: Sequence[Dict[str, Any]]) -> List[Dict[str, Any]]:
    slides: List[Dict[str, Any]] = []
    for title in _normalize_str_list(requirement.get("knowledge_points")):
        short_title = _shorten_slide_title(title, [], "")
        slides.append(
            {
                "slide_type": "content",
                "title": short_title,
                "toc_title": _build_toc_title(short_title, title),
                "goal": "",
                "bullets": _prepare_display_bullets(_clip_bullets(_normalize_str_list(requirement.get("key_points")) or [title, "Use a representative example"], minimum=2), short_title)[:4],
                "visual_suggestion": short_title,
                "source_refs": [],
            }
        )

    for item in docx_outline:
        if not isinstance(item, dict):
            continue
        title = str(item.get("section_title") or "").strip()
        bullets = _clip_bullets(item.get("bullets"), minimum=2)
        if not title and not bullets:
            continue
        short_title = _shorten_slide_title(title, bullets, str(item.get("section_goal") or "").strip())
        slides.append(
            {
                "slide_type": "content",
                "title": short_title or "Teaching Content",
                "toc_title": _build_toc_title(short_title or "Teaching Content", title),
                "goal": str(item.get("section_goal") or "").strip(),
                "bullets": _prepare_display_bullets(bullets if bullets else ["Advance the lesson through guided activity"], short_title or title)[:4],
                "visual_suggestion": short_title or title,
                "source_refs": [str(value).strip() for value in item.get("source_refs", []) if str(value).strip()],
            }
        )

    return slides


def _build_summary_slide(
    summary_candidates: Sequence[Dict[str, Any]],
    requirement: Dict[str, Any],
    content_slides: Sequence[Dict[str, Any]],
) -> Dict[str, Any]:
    if summary_candidates:
        candidate = _normalize_outline_slide(summary_candidates[0])
        if candidate["title"] or candidate["bullets"]:
            return {"title": candidate["title"] or "Lesson Summary", "bullets": candidate["bullets"][:4] or ["Review the key ideas from this lesson"]}

    bullets = _clip_bullets(requirement.get("teaching_goals"), minimum=2)
    if len(bullets) < 4:
        bullets.extend(_clip_bullets(requirement.get("key_points"), minimum=0))
    if len(bullets) < 4:
        bullets.extend([item.get("title", "") for item in content_slides[:2] if item.get("title")])
    bullets = [item for item in bullets if item][:4]
    if not bullets:
        bullets = ["Review the key concepts", "Summarize the knowledge structure", "Connect activity and practice", "Prepare for next-step learning"]
    return {"title": "Lesson Summary", "bullets": bullets}


def _infer_theme_from_spec(spec: Dict[str, Any]) -> str:
    requirement = spec.get("requirement_summary") if isinstance(spec.get("requirement_summary"), dict) else {}
    style = requirement.get("style") if isinstance(requirement.get("style"), dict) else {}
    teaching_style = str(style.get("teaching_style") or "").strip().lower()
    joined = f" {teaching_style} "
    if any(token in joined for token in ["explore", "experiment", "technology", "tech", "实验", "信息化", "技术"]):
        return "tech"
    if any(token in joined for token in ["interactive", "vivid", "game", "scenario", "互动", "活泼", "游戏", "情境"]):
        return "vivid"
    return "clean"


def _extract_slide_keywords(slide: Dict[str, Any]) -> List[str]:
    values = [
        str(slide.get("title") or "").strip(),
        str(slide.get("goal") or "").strip(),
        str(slide.get("visual_suggestion") or "").strip(),
    ]
    values.extend([str(item).strip() for item in slide.get("bullets", []) if str(item).strip()])
    values.extend([str(item).strip() for item in slide.get("source_refs", []) if str(item).strip()])
    seen = set()
    tokens: List[str] = []
    for value in values:
        for token in _split_words(value):
            if token and token not in seen:
                seen.add(token)
                tokens.append(token)
    return tokens


def _split_words(value: str) -> List[str]:
    words = [
        item.strip().lower()
        for item in value.replace("，", " ").replace("。", " ").replace("：", " ").replace(",", " ").split()
    ]
    tokens: List[str] = []
    for word in words:
        if not word:
            continue
        tokens.append(word)
        if len(word) >= 4:
            tokens.append(word[:4])
    if not tokens and value.strip():
        tokens.append(value.strip().lower())
    return tokens


def _score_keyword_match(text: str, keywords: Sequence[str]) -> int:
    haystack = str(text or "").strip().lower()
    if not haystack:
        return 0
    return sum(5 for keyword in keywords[:8] if keyword and keyword in haystack)


def _slide_ref_matches(slide: Dict[str, Any], mapping: Dict[str, Any]) -> bool:
    slide_refs = [str(value).strip() for value in slide.get("source_refs", []) if str(value).strip()]
    if not slide_refs:
        return False
    file_name = str(mapping.get("file_name") or "").strip()
    knowledge_point = str(mapping.get("knowledge_point") or "").strip()
    return any(ref in file_name or ref in knowledge_point for ref in slide_refs)


def _build_gallery_queries(slide: Dict[str, Any], spec: Dict[str, Any]) -> List[str]:
    requirement = spec.get("requirement_summary") if isinstance(spec.get("requirement_summary"), dict) else {}
    knowledge_points = _normalize_str_list(requirement.get("knowledge_points"))
    base_query = " ".join(
        [
            item
            for item in [
                str(slide.get("title") or "").strip(),
                str(slide.get("visual_suggestion") or "").strip(),
                str(requirement.get("grade_subject") or "").strip(),
                " ".join(knowledge_points[:2]),
            ]
            if item
        ]
    ).strip()
    queries = [base_query] if base_query else []
    title = str(slide.get("title") or "").strip()
    if title and title not in queries:
        queries.append(title)
    if knowledge_points:
        kp_query = " ".join(knowledge_points[:2]).strip()
        if kp_query and kp_query not in queries:
            queries.append(kp_query)
    return queries or ["education classroom learning"]


def _download_gallery_image(image_url: str, photo_id: str, query: str, course_id: int, upload_root: str) -> Optional[str]:
    assets_dir = os.path.join(upload_root, "materials", str(course_id), ".generated_assets")
    os.makedirs(assets_dir, exist_ok=True)
    filename = f"pexels_{photo_id}_{hashlib.sha1(query.encode('utf-8')).hexdigest()[:8]}.jpg"
    abs_path = os.path.join(assets_dir, filename)
    if os.path.exists(abs_path):
        return abs_path

    response = requests.get(image_url, timeout=15)
    response.raise_for_status()
    with open(abs_path, "wb") as file_obj:
        file_obj.write(response.content)
    return abs_path


def _render_cover_slide(prs: Presentation, cover: Dict[str, Any], theme: Dict[str, Any], spec: Dict[str, Any], layout_map: Optional[Dict[str, Any]] = None) -> None:
    layout_map = layout_map or {}
    layout = _get_layout(prs, layout_map, "cover")
    slide = prs.slides.add_slide(layout)
    _paint_slide_background(slide, theme, "cover")
    _fill_placeholder_or_textbox(slide, 0, str(cover.get("title") or "AI Lesson Deck"), 0.9, 1.7, 11.0, 1.4, theme["cover_title_size"], theme["text"], theme["font_family"], True)
    _fill_placeholder_or_textbox(slide, 1, str(cover.get("subtitle") or "易度新星 EduNova generated presentation"), 0.95, 3.0, 8.8, 0.7, 18, theme["muted"], theme["font_family"])


def _render_toc_slide(prs: Presentation, toc: Dict[str, Any], theme: Dict[str, Any], layout_map: Optional[Dict[str, Any]] = None) -> None:
    layout_map = layout_map or {}
    layout = _get_layout(prs, layout_map, "toc")
    slide = prs.slides.add_slide(layout)
    _paint_slide_background(slide, theme, "toc")
    _add_page_title(slide, "Agenda", theme)
    items = toc.get("items") if isinstance(toc.get("items"), list) else []
    columns = [items[:4], items[4:8]]

    for column_index, column_items in enumerate(columns):
        top = 1.8
        left = 0.9 if column_index == 0 else 6.9
        for item_index, text in enumerate(column_items, start=1 + column_index * 4):
            _add_outline_item(slide, str(text or f"Topic {item_index}"), item_index, theme, left, top)
            top += 1.15
    _add_footer(slide, 2, theme)


def _render_content_slide(
    prs: Presentation,
    slide_data: Dict[str, Any],
    theme: Dict[str, Any],
    image: Optional[Dict[str, Any]],
    layout_name: str,
    page_number: int,
    game_entry: Optional[Dict[str, Any]] = None,
    layout_map: Optional[Dict[str, Any]] = None,
) -> None:
    layout_map = layout_map or {}
    if layout_name == "text-focus":
        layout_key = "content_text"
    elif layout_name == "image-left":
        layout_key = "content_img_left"
    else:
        layout_key = "content_img_right"
    layout = _get_layout(prs, layout_map, layout_key)
    slide = prs.slides.add_slide(layout)
    _paint_slide_background(slide, theme, "content")
    title_text = str(slide_data.get("title") or "Content Slide")
    title_bottom = _add_page_title(slide, title_text, theme)
    content_top = max(1.55, title_bottom + 0.15)

    if layout_name == "text-focus":
        text_left, text_top, text_width, text_height = 0.95, content_top, 11.35, max(3.9, 6.7 - content_top)
    elif layout_name == "image-left":
        _add_image_block(slide, image, 0.7, content_top, 5.0, max(3.8, 6.45 - content_top), theme)
        text_left, text_top, text_width, text_height = 6.0, content_top, 6.35, max(3.9, 6.6 - content_top)
    else:
        _add_image_block(slide, image, 7.6, content_top, 5.0, max(3.8, 6.45 - content_top), theme)
        text_left, text_top, text_width, text_height = 0.8, content_top, 6.2, max(3.9, 6.6 - content_top)

    _add_bullet_box(slide, [str(item).strip() for item in slide_data.get("bullets", []) if str(item).strip()][:4], text_left, text_top, text_width, text_height, theme)
    _add_homework_game_link(slide, slide_data, theme, game_entry)
    _add_footer(slide, page_number, theme)


def _is_homework_slide(slide_data: Dict[str, Any]) -> bool:
    title_text = str(slide_data.get("title") or "").strip().lower()
    bullets_text = " ".join(str(item).strip().lower() for item in slide_data.get("bullets", []) if str(item).strip())
    combined = f"{title_text} {bullets_text}"
    keywords = ("课后", "作业", "homework", "after-class", "练习任务")
    return any(keyword in combined for keyword in keywords)


def _add_homework_game_link(
    slide: Any,
    slide_data: Dict[str, Any],
    theme: Dict[str, Any],
    game_entry: Optional[Dict[str, Any]],
) -> None:
    if not game_entry or not _is_homework_slide(slide_data):
        return
    game_url = str(game_entry.get("download_url") or "").strip()
    if not game_url:
        return

    _add_textbox(slide, "互动小游戏链接", 0.95, 6.12, 2.0, 0.28, 11, theme["muted"], theme["font_family"], True)
    _add_hyperlink_textbox(
        slide,
        game_url,
        0.95,
        6.36,
        11.1,
        0.36,
        12,
        theme["accent"],
        theme["font_family"],
        hyperlink=game_url,
    )


def _render_summary_slide(prs: Presentation, summary: Dict[str, Any], theme: Dict[str, Any], page_number: int, layout_map: Optional[Dict[str, Any]] = None) -> None:
    layout_map = layout_map or {}
    layout = _get_layout(prs, layout_map, "summary")
    slide = prs.slides.add_slide(layout)
    _paint_slide_background(slide, theme, "summary")
    _add_page_title(slide, str(summary.get("title") or "Lesson Summary"), theme)
    _add_summary_cards(slide, [str(item).strip() for item in summary.get("bullets", []) if str(item).strip()][:4], theme)
    _add_footer(slide, page_number, theme)


def _render_ending_slide(prs: Presentation, ending: Dict[str, Any], theme: Dict[str, Any], page_number: int, layout_map: Optional[Dict[str, Any]] = None) -> None:
    layout_map = layout_map or {}
    layout = _get_layout(prs, layout_map, "ending")
    slide = prs.slides.add_slide(layout)
    _paint_slide_background(slide, theme, "ending")
    _add_textbox(slide, str(ending.get("title") or "Thank You"), 1.2, 2.2, 10.8, 1.1, theme["cover_title_size"], theme["text"], theme["font_family"], True, PP_ALIGN.CENTER)
    _add_textbox(slide, str(ending.get("subtitle") or ""), 1.6, 3.45, 10.0, 0.7, 18, theme["muted"], theme["font_family"], False, PP_ALIGN.CENTER)
    _add_footer(slide, page_number, theme)


def _render_game_entry_slide(prs: Presentation, game_entry: Dict[str, Any], theme: Dict[str, Any], page_number: int, layout_map: Optional[Dict[str, Any]] = None) -> None:
    layout_map = layout_map or {}
    layout = _get_layout(prs, layout_map, "ending")
    slide = prs.slides.add_slide(layout)
    _paint_slide_background(slide, theme, "ending")
    _add_page_title(slide, str(game_entry.get("title") or "小游戏入口页"), theme)
    _add_textbox(
        slide,
        str(game_entry.get("subtitle") or "扫码或点击链接下载小游戏"),
        0.9,
        1.5,
        7.2,
        0.5,
        16,
        theme["muted"],
        theme["font_family"],
    )

    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.0), Inches(7.2), Inches(3.9))
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(theme["summary_bg"])
    card.line.color.rgb = _rgb(theme["accent_soft"])
    _set_shape_text(card, "", theme["font_family"], 14, theme["text"], False, PP_ALIGN.LEFT)

    _add_textbox(slide, "文件", 1.2, 2.35, 1.0, 0.3, 12, theme["muted"], theme["font_family"], True)
    _add_textbox(slide, str(game_entry.get("material_title") or ""), 1.2, 2.63, 5.8, 0.5, 18, theme["text"], theme["font_family"], True)
    _add_textbox(slide, "下载链接", 1.2, 3.35, 1.3, 0.3, 12, theme["muted"], theme["font_family"], True)
    _add_hyperlink_textbox(
        slide,
        str(game_entry.get("download_url") or ""),
        1.2,
        3.66,
        5.9,
        1.2,
        14,
        theme["accent"],
        theme["font_family"],
        hyperlink=str(game_entry.get("download_url") or "").strip() or None,
    )
    _add_textbox(slide, "教师提示", 1.2, 5.05, 1.3, 0.3, 12, theme["muted"], theme["font_family"], True)
    _add_textbox(slide, str(game_entry.get("teacher_tip") or ""), 1.2, 5.35, 5.9, 0.55, 12, theme["text"], theme["font_family"])

    qr_path = str(game_entry.get("qr_image_path") or "").strip()
    qr_frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.65), Inches(2.0), Inches(3.0), Inches(3.0))
    qr_frame.fill.solid()
    qr_frame.fill.fore_color.rgb = _rgb((255, 255, 255))
    qr_frame.line.color.rgb = _rgb(theme["accent_soft"])
    if qr_path and os.path.exists(qr_path):
        _add_picture_fit(slide, qr_path, 8.82, 2.17, 2.66, 2.66)
    else:
        _set_shape_text(qr_frame, "QR unavailable", theme["font_family"], 14, theme["muted"], False, PP_ALIGN.CENTER)

    _add_textbox(slide, "扫码下载小游戏", 8.7, 5.2, 2.9, 0.35, 13, theme["text"], theme["font_family"], True, PP_ALIGN.CENTER)
    _add_textbox(slide, "如扫码受限，可直接点击左侧链接。", 8.55, 5.56, 3.2, 0.38, 10, theme["muted"], theme["font_family"], False, PP_ALIGN.CENTER)
    _add_footer(slide, page_number, theme)


def _fill_placeholder_or_textbox(
    slide: Any,
    ph_idx: int,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int,
    color: Tuple[int, int, int],
    font_family: str,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """Try to fill a layout placeholder by *ph_idx*; fall back to a free textbox."""
    placeholder = None
    try:
        placeholder = slide.placeholders[ph_idx]
    except (KeyError, IndexError):
        pass

    if placeholder is not None:
        text_frame = placeholder.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        paragraph = text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = alignment
        paragraph.font.name = font_family
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = _rgb(color)
    else:
        _add_textbox(slide, text, left, top, width, height, font_size, color, font_family, bold, alignment)


def _paint_slide_background(slide: Any, theme: Dict[str, Any], variant: str) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = _rgb(theme["bg"])

    top_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.18))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = _rgb(theme["accent"])
    top_bar.line.fill.background()

    if variant != "cover":
        accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, SLIDE_WIDTH - Inches(1.55), Inches(0.35), Inches(1.05), Inches(0.2))
        accent.fill.solid()
        accent.fill.fore_color.rgb = _rgb(theme["accent"])
        accent.fill.transparency = 0.15
        accent.line.fill.background()

    if variant in {"cover", "ending"}:
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, SLIDE_WIDTH - Inches(2.2), SLIDE_HEIGHT - Inches(2.3), Inches(2.7), Inches(2.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = _rgb(theme["accent_soft"])
        circle.fill.transparency = 0.35
        circle.line.fill.background()


def _add_page_title(slide: Any, text: str, theme: Dict[str, Any]) -> float:
    cleaned = str(text or "").strip()
    font_size = theme["title_size"]
    height = 0.7
    if len(cleaned) > 24:
        font_size = max(20, theme["title_size"] - 4)
        height = 0.95
    if len(cleaned) > 38:
        font_size = max(18, theme["title_size"] - 7)
        height = 1.2
    _add_textbox(slide, cleaned, 0.9, 0.55, 8.9, height, font_size, theme["text"], theme["font_family"], True)
    return 0.55 + height


def _add_outline_item(slide: Any, text: str, number: int, theme: Dict[str, Any], left: float, top: float) -> None:
    number_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(0.78), Inches(0.54))
    number_shape.fill.solid()
    number_shape.fill.fore_color.rgb = _rgb(theme["accent"])
    number_shape.line.fill.background()
    _set_shape_text(number_shape, f"{number:02d}", theme["font_family"], 14, (255, 255, 255), True, PP_ALIGN.CENTER)
    display_text = _build_toc_title(_shorten_slide_title(text, [], ""), text)
    _add_textbox(slide, display_text, left + 1.0, top - 0.02, 4.5, 0.65, 18, theme["text"], theme["font_family"])


def _add_bullet_box(slide: Any, bullets: Sequence[str], left: float, top: float, width: float, height: float, theme: Dict[str, Any]) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Pt(8)
    text_frame.margin_right = Pt(8)
    text_frame.margin_top = Pt(6)
    text_frame.margin_bottom = Pt(6)
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = f"• {bullet}"
        paragraph.level = 0
        paragraph.space_after = Pt(10)
        paragraph.font.name = theme["font_family"]
        paragraph.font.size = Pt(theme["body_size"])
        paragraph.font.color.rgb = _rgb(theme["text"])


def _add_summary_cards(slide: Any, bullets: Sequence[str], theme: Dict[str, Any]) -> None:
    bullets = list(bullets) or ["Review the key concept", "Summarize the method", "Link practice and transfer", "Clarify the next step"]
    positions = [(0.9, 1.85), (6.8, 1.85), (0.9, 4.1), (6.8, 4.1)]
    for index, bullet in enumerate(bullets[:4]):
        left, top = positions[index]
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.55), Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(theme["summary_bg"])
        card.line.color.rgb = _rgb(theme["accent_soft"])
        _set_shape_text(card, bullet, theme["font_family"], theme["body_size"], theme["text"], False, PP_ALIGN.LEFT)


def _add_image_block(slide: Any, image: Optional[Dict[str, Any]], left: float, top: float, width: float, height: float, theme: Dict[str, Any]) -> None:
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    frame.fill.solid()
    frame.fill.fore_color.rgb = _rgb(theme["accent_soft"])
    frame.line.fill.background()

    if not image:
        _set_shape_text(frame, "No image", theme["font_family"], 16, theme["muted"], False, PP_ALIGN.CENTER)
        return

    image_path = str(image.get("path_or_url") or "").strip()
    if not image_path or not os.path.exists(image_path):
        _set_shape_text(frame, "Image unavailable", theme["font_family"], 16, theme["muted"], False, PP_ALIGN.CENTER)
        return

    frame.fill.background()
    _add_picture_fit(slide, image_path, left, top, width, height)
    caption = str(image.get("caption") or "").strip()
    if caption:
        _add_caption(slide, caption, theme, left, top + height + 0.08, width)


def _add_picture_fit(slide: Any, image_path: str, left: float, top: float, width: float, height: float) -> None:
    with Image.open(image_path) as image:
        image_width, image_height = image.size
    if not image_width or not image_height:
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width), height=Inches(height))
        return

    target_ratio = width / height
    image_ratio = image_width / image_height
    if image_ratio > target_ratio:
        picture_height = height
        picture_width = height * image_ratio
    else:
        picture_width = width
        picture_height = width / image_ratio

    left_offset = left - max(0, (picture_width - width) / 2)
    top_offset = top - max(0, (picture_height - height) / 2)
    slide.shapes.add_picture(image_path, Inches(left_offset), Inches(top_offset), width=Inches(picture_width), height=Inches(picture_height))


def _add_caption(slide: Any, text: str, theme: Dict[str, Any], left: float, top: float, width: float) -> None:
    _add_textbox(slide, text, left, top, width, 0.4, 11, theme["muted"], theme["font_family"])


def _add_footer(slide: Any, page_number: int, theme: Dict[str, Any]) -> None:
    footer_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.85), Inches(7.0), Inches(11.65), Inches(0.04))
    footer_line.fill.solid()
    footer_line.fill.fore_color.rgb = _rgb(theme["accent_soft"])
    footer_line.line.fill.background()
    _add_textbox(slide, f"{page_number:02d}", 12.15, 6.78, 0.5, 0.25, theme["small_size"], theme["footer"], theme["font_family"], False, PP_ALIGN.RIGHT)


def _add_band(slide: Any, color: Tuple[int, int, int], y: float, height: float) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(y), Inches(4.7), Inches(height))
    band.fill.solid()
    band.fill.fore_color.rgb = _rgb(color)
    band.fill.transparency = 0.12
    band.line.fill.background()


def _add_label(slide: Any, text: str, theme: Dict[str, Any], left: float, top: float, width: float) -> None:
    label = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.42))
    label.fill.solid()
    label.fill.fore_color.rgb = _rgb(theme["accent"])
    label.line.fill.background()
    _set_shape_text(label, text, theme["font_family"], 11, (255, 255, 255), True, PP_ALIGN.CENTER)


def _add_textbox(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int,
    color: Tuple[int, int, int],
    font_family: str,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = alignment
    paragraph.font.name = font_family
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = _rgb(color)
    return textbox


def _add_hyperlink_textbox(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int,
    color: Tuple[int, int, int],
    font_family: str,
    hyperlink: Optional[str] = None,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = alignment
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_family
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    if hyperlink:
        run.hyperlink.address = hyperlink


def _set_shape_text(shape: Any, text: str, font_family: str, font_size: int, color: Tuple[int, int, int], bold: bool, alignment: PP_ALIGN) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = alignment
    paragraph.font.name = font_family
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = _rgb(color)


def _build_candidate(source_type: str, path_or_url: str, score: int, caption: str, source_ref: str) -> Dict[str, Any]:
    return {"source_type": source_type, "path_or_url": path_or_url, "score": score, "caption": caption, "source_ref": source_ref}


def _clip_bullets(raw_items: Any, minimum: int = 0) -> List[str]:
    items = _normalize_str_list(raw_items)
    return items[:5] if len(items) >= minimum else items


def _prepare_display_bullets(bullets: Sequence[str], title: str) -> List[str]:
    results: List[str] = []
    seen = set()
    for bullet in bullets:
        text = _compress_ppt_bullet(str(bullet or "").strip(), title)
        if text and text not in seen:
            seen.add(text)
            results.append(text)
    return results


def _expand_slide_bullets(bullets: Sequence[str], title: str, goal: str) -> List[str]:
    topic = str(title or goal or "本页内容").strip()
    expanded: List[str] = []
    seen = set()
    for bullet in bullets:
        text = str(bullet or "").strip()
        if not text:
            continue
        if len(text) < 16:
            text = f"{text}，建议结合“{topic}”补充一个课堂例子或提问点，帮助学生快速理解。"
        if text not in seen:
            seen.add(text)
            expanded.append(text)
    return expanded


def _compress_ppt_bullet(text: str, title: str) -> str:
    cleaned = re.sub(r"\s+", " ", str(text or "").strip())
    if not cleaned:
        return ""
    if len(cleaned) <= 34:
        return cleaned
    clauses = [part.strip(" ，；。,:：") for part in re.split(r"[；。]", cleaned) if part.strip(" ，；。,:：")]
    if clauses:
        first = clauses[0]
        if len(first) <= 34:
            return first
        comma_parts = [part.strip(" ，；。,:：") for part in re.split(r"[，,:：]", first) if part.strip(" ，；。,:：")]
        if comma_parts:
            shortened = "，".join(comma_parts[:2]).strip(" ，；。,:：")
            if shortened:
                return shortened[:34] + ("…" if len(shortened) > 34 else "")
    if "：" in cleaned:
        prefix, suffix = cleaned.split("：", 1)
        prefix = prefix.strip()
        suffix = suffix.strip()
        if prefix and suffix:
            tail = re.split(r"[，；。,:：]", suffix)[0].strip()
            candidate = f"{prefix}：{tail}".strip(" ：")
            if candidate:
                return candidate[:34] + ("…" if len(candidate) > 34 else "")
    return cleaned[:34] + "…"


def _shorten_slide_title(text: str, bullets: Sequence[str], goal: str) -> str:
    cleaned = re.sub(r"\s+", " ", str(text or "").strip())
    if not cleaned:
        fallback = str(goal or "").strip()
        return fallback[:20] + ("…" if len(fallback) > 20 else "") if fallback else "课堂内容"
    if len(cleaned) <= 24:
        return cleaned
    if "：" in cleaned:
        prefix, suffix = cleaned.split("：", 1)
        prefix = prefix.strip()
        suffix = suffix.strip()
        head = re.split(r"[，；。,:：]", suffix)[0].strip()
        if prefix and head:
            candidate = f"{prefix}：{head}"
            return candidate[:24] + ("…" if len(candidate) > 24 else "")
    short_clause = re.split(r"[，；。,:：]", cleaned)[0].strip()
    if short_clause:
        return short_clause[:24] + ("…" if len(short_clause) > 24 else "")
    fallback = str(goal or (bullets[0] if bullets else cleaned)).strip()
    return fallback[:24] + ("…" if len(fallback) > 24 else "")


def _extract_title_detail_as_bullet(raw_title: str, short_title: str) -> str:
    raw = str(raw_title or "").strip()
    short = str(short_title or "").strip().rstrip("…")
    if not raw or len(raw) <= 24:
        return ""
    if short and raw.startswith(short):
        remainder = raw[len(short):].strip("：:，；。 ")
        return _compress_ppt_bullet(remainder, short) if remainder else ""
    if "：" in raw:
        _, suffix = raw.split("：", 1)
        return _compress_ppt_bullet(suffix, short or raw)
    return _compress_ppt_bullet(raw, short or raw)


def _build_toc_title(short_title: str, raw_title: str) -> str:
    base = str(short_title or "").strip() or str(raw_title or "").strip() or "课堂内容"
    return base[:22] + ("…" if len(base) > 22 else "")


def _split_dense_content_slides(slides: Sequence[Dict[str, Any]], desired_count: int) -> List[Dict[str, Any]]:
    result = [dict(item) for item in slides if isinstance(item, dict)]
    while len(result) < desired_count:
        split_index = -1
        for index, slide in enumerate(result):
            bullets = [str(item).strip() for item in slide.get("bullets", []) if str(item).strip()]
            total_chars = sum(len(item) for item in bullets)
            if len(bullets) >= 4 and total_chars >= 72:
                split_index = index
                break
        if split_index < 0:
            break
        target = dict(result[split_index])
        bullets = [str(item).strip() for item in target.get("bullets", []) if str(item).strip()]
        midpoint = max(2, len(bullets) // 2)
        first_half = bullets[:midpoint]
        second_half = bullets[midpoint:]
        if len(second_half) < 2:
            break
        target["bullets"] = first_half
        continuation = dict(target)
        continuation["title"] = _shorten_slide_title(f"{target.get('title') or '课堂内容'}（续）", second_half, str(target.get("goal") or ""))
        continuation["toc_title"] = _build_toc_title(continuation["title"], continuation["title"])
        continuation["bullets"] = second_half
        result[split_index] = target
        result.insert(split_index + 1, continuation)
    return result


def _normalize_str_list(value: Any) -> List[str]:
    if not isinstance(value, list):
        return []
    result: List[str] = []
    seen = set()
    for item in value:
        text = str(item or "").strip()
        if text and text not in seen:
            seen.add(text)
            result.append(text)
    return result


def _resolve_upload_path(upload_root: str, relative_path: Any) -> Optional[str]:
    if not isinstance(relative_path, str):
        return None
    normalized = os.path.normpath(relative_path.lstrip("/\\")).replace("\\", "/")
    if not normalized or normalized.startswith(".."):
        return None
    abs_path = os.path.abspath(os.path.join(upload_root, normalized))
    if not abs_path.startswith(os.path.abspath(upload_root)) or not os.path.exists(abs_path):
        return None
    return abs_path


def _calculate_file_hash(file_path: str) -> Optional[str]:
    try:
        digest = hashlib.sha256()
        with open(file_path, "rb") as file_obj:
            for chunk in iter(lambda: file_obj.read(4096), b""):
                digest.update(chunk)
        return digest.hexdigest()
    except Exception:
        return None


def _material_upload_to_abs_path(file_path: str) -> Optional[str]:
    if not isinstance(file_path, str) or not file_path:
        return None
    project_backend_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    return os.path.join(project_backend_dir, file_path.lstrip("/"))


def _build_qr_code_image(download_url: str, assets_dir: str) -> str:
    digest = hashlib.sha1(download_url.encode("utf-8")).hexdigest()[:12]
    qr_path = os.path.join(assets_dir, f"game_entry_qr_{digest}.png")
    if os.path.exists(qr_path):
        return qr_path

    qr = qrcode.QRCode(
        version=None,
        error_correction=qrcode.constants.ERROR_CORRECT_M,
        box_size=8,
        border=2,
    )
    qr.add_data(download_url)
    qr.make(fit=True)
    image = qr.make_image(fill_color="black", back_color="white")
    image.save(qr_path)
    return qr_path


def _rgb(color: Tuple[int, int, int]) -> RGBColor:
    return RGBColor(color[0], color[1], color[2])


def _landscape_bonus(path_or_url: Any) -> int:
    if not isinstance(path_or_url, str) or not os.path.exists(path_or_url):
        return 0
    try:
        with Image.open(path_or_url) as image:
            width, height = image.size
        return 2 if width >= height else 0
    except Exception:
        return 0
