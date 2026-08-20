import os
from typing import Any, Dict, List, Optional

from . import (
    PARSER_VERSION,
    ParseResult,
    build_text_chunks,
    calculate_file_hash,
    clip_text,
    ensure_derived_dir,
    guess_extension_from_content_type,
    load_cached_parse_result,
    run_soffice_convert,
    to_rel_upload_path,
    write_cached_parse_result,
)
from .ocr import run_remote_ocr


def parse_ppt(
    file_path: str,
    *,
    upload_root: Optional[str] = None,
    owner_id: Any = "system",
    file_hash: Optional[str] = None,
    api_key: Optional[str] = None,
    api_base: Optional[str] = None,
    parse_mode: str = "ppt",
) -> ParseResult:
    resolved_hash = file_hash or calculate_file_hash(file_path)
    cached = load_cached_parse_result(upload_root, resolved_hash, parse_mode)
    if cached:
        return cached

    derived_dir = ensure_derived_dir(upload_root, owner_id, resolved_hash, namespace="ppt") if upload_root else None
    source_path = file_path
    if file_path.lower().endswith(".ppt"):
        if not derived_dir:
            raise RuntimeError("旧版 PPT 解析需要可用的上传根目录")
        source_path = _convert_legacy_ppt_to_pptx(file_path, derived_dir)

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PPT parsing") from exc

    presentation = Presentation(source_path)
    slides: List[Dict[str, Any]] = []
    raw_slide_texts: List[str] = []
    has_ocr = False

    for slide_index, slide in enumerate(presentation.slides, start=1):
        parsed_slide = _parse_slide(
            slide,
            slide_index=slide_index,
            mso_shape_type_group=MSO_SHAPE_TYPE.GROUP,
            mso_shape_type_picture=MSO_SHAPE_TYPE.PICTURE,
            derived_dir=derived_dir,
            upload_root=upload_root,
            api_key=api_key,
            api_base=api_base,
        )
        if parsed_slide.get("ocr_text"):
            has_ocr = True
        slides.append(parsed_slide)
        slide_text = _build_slide_text(parsed_slide)
        parsed_slide["speaker_notes_weighted_text"] = slide_text
        if slide_text:
            raw_slide_texts.append(slide_text)

    raw_text = "\n\n".join(raw_slide_texts).strip()
    chunks = _build_slide_chunks(slides)
    result: ParseResult = {
        "raw_text": raw_text,
        "summary": _summarize_slides(slides),
        "chunks": chunks,
        "structure": {"slides": slides},
        "assets": {"slides": slides, "slide_count": len(slides)},
        "meta": {
            "parser_version": PARSER_VERSION,
            "file_type": "pptx" if source_path.lower().endswith(".pptx") else "ppt",
            "source_path": file_path,
            "has_ocr": has_ocr,
            "scan_detected": False,
        },
    }
    write_cached_parse_result(upload_root, resolved_hash, parse_mode, result)
    return result


def _convert_legacy_ppt_to_pptx(full_path: str, output_dir: str) -> str:
    abs_path = os.path.abspath(full_path)
    expected_output = os.path.join(output_dir, os.path.splitext(os.path.basename(abs_path))[0] + ".pptx")
    if not os.path.exists(expected_output):
        run_soffice_convert(
            source_path=abs_path,
            output_dir=output_dir,
            convert_to="pptx",
            missing_message="未安装 LibreOffice，暂时不能解析ppt格式",
        )
    if not os.path.exists(expected_output):
        raise RuntimeError("PPT 已调用 LibreOffice 转换，但未生成 .pptx 文件")
    return expected_output


def _parse_slide(
    slide: Any,
    *,
    slide_index: int,
    mso_shape_type_group: Any,
    mso_shape_type_picture: Any,
    derived_dir: Optional[str],
    upload_root: Optional[str],
    api_key: Optional[str],
    api_base: Optional[str],
) -> Dict[str, Any]:
    slide_data: Dict[str, Any] = {
        "slide_index": slide_index,
        "title": "",
        "text": "",
        "notes": _extract_slide_notes(slide),
        "tables": [],
        "images": [],
        "ocr_text": "",
        "speaker_notes_weighted_text": "",
    }
    text_parts: List[str] = []
    table_texts: List[str] = []
    image_texts: List[str] = []

    try:
        if slide.shapes.title and slide.shapes.title.text:
            slide_data["title"] = str(slide.shapes.title.text).strip()
    except Exception:
        slide_data["title"] = ""

    for shape in slide.shapes:
        _collect_shape_content(
            shape,
            mso_shape_type_group=mso_shape_type_group,
            mso_shape_type_picture=mso_shape_type_picture,
            slide_index=slide_index,
            text_parts=text_parts,
            table_texts=table_texts,
            images=slide_data["images"],
            image_texts=image_texts,
            derived_dir=derived_dir,
            upload_root=upload_root,
            api_key=api_key,
            api_base=api_base,
        )

    slide_data["tables"] = [{"table_index": index, "text": text} for index, text in enumerate(table_texts, start=1)]
    slide_data["text"] = "\n".join(part for part in text_parts if part).strip()
    slide_data["ocr_text"] = "\n".join(part for part in image_texts if part).strip()
    return slide_data


def _collect_shape_content(
    shape: Any,
    *,
    mso_shape_type_group: Any,
    mso_shape_type_picture: Any,
    slide_index: int,
    text_parts: List[str],
    table_texts: List[str],
    images: List[Dict[str, Any]],
    image_texts: List[str],
    derived_dir: Optional[str],
    upload_root: Optional[str],
    api_key: Optional[str],
    api_base: Optional[str],
) -> None:
    shape_type = getattr(shape, "shape_type", None)

    if getattr(shape, "has_text_frame", False):
        text = "\n".join(
            paragraph.text.strip()
            for paragraph in shape.text_frame.paragraphs
            if paragraph.text and paragraph.text.strip()
        ).strip()
        if text:
            text_parts.append(text)

    if getattr(shape, "has_table", False):
        table_text = _extract_table_text(shape.table)
        if table_text:
            table_texts.append(table_text)

    if getattr(shape, "has_chart", False):
        chart_text = _extract_chart_text(shape)
        if chart_text:
            text_parts.append(chart_text)

    if shape_type == mso_shape_type_group:
        for child_shape in shape.shapes:
            _collect_shape_content(
                child_shape,
                mso_shape_type_group=mso_shape_type_group,
                mso_shape_type_picture=mso_shape_type_picture,
                slide_index=slide_index,
                text_parts=text_parts,
                table_texts=table_texts,
                images=images,
                image_texts=image_texts,
                derived_dir=derived_dir,
                upload_root=upload_root,
                api_key=api_key,
                api_base=api_base,
            )
        return

    if shape_type == mso_shape_type_picture and derived_dir and upload_root:
        try:
            image = shape.image
            extension = guess_extension_from_content_type(getattr(image, "content_type", ""), fallback=".png")
            image_abs_path = os.path.join(derived_dir, f"slide_{slide_index:03d}_image_{len(images) + 1:02d}{extension}")
            with open(image_abs_path, "wb") as file_obj:
                file_obj.write(image.blob)
            ocr_text = ""
            if api_key and api_base:
                ocr_result = run_remote_ocr(image_abs_path, api_key, api_base)
                ocr_text = str(ocr_result.get("raw_text") or "").strip()
            images.append({"image_path": to_rel_upload_path(image_abs_path, upload_root), "ocr_text": ocr_text})
            if ocr_text:
                image_texts.append(ocr_text)
        except Exception:
            return


def _extract_table_text(table: Any) -> str:
    rows: List[str] = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cell_text = " ".join(
                paragraph.text.strip()
                for paragraph in cell.text_frame.paragraphs
                if paragraph.text and paragraph.text.strip()
            ).strip()
            cells.append(cell_text)
        if any(cells):
            rows.append(" | ".join(cell for cell in cells if cell))
    return "\n".join(rows).strip()


def _extract_chart_text(shape: Any) -> str:
    parts: List[str] = []
    try:
        chart = shape.chart
    except Exception:
        return ""
    try:
        if chart.has_title and chart.chart_title and chart.chart_title.text_frame:
            title = chart.chart_title.text_frame.text.strip()
            if title:
                parts.append(title)
    except Exception:
        pass
    try:
        series_names = [str(series.name).strip() for series in chart.series if str(series.name).strip()]
        if series_names:
            parts.append("系列: " + "、".join(series_names))
    except Exception:
        pass
    return "\n".join(parts).strip()


def _extract_slide_notes(slide: Any) -> str:
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return ""
    try:
        notes_frame = notes_slide.notes_text_frame
        if notes_frame:
            return "\n".join(
                paragraph.text.strip()
                for paragraph in notes_frame.paragraphs
                if paragraph.text and paragraph.text.strip()
            ).strip()
    except Exception:
        return ""
    return ""


def _build_slide_text(slide: Dict[str, Any]) -> str:
    ordered_parts = [
        str(slide.get("title") or "").strip(),
        str(slide.get("text") or "").strip(),
        "\n".join(str(item.get("text") or "").strip() for item in slide.get("tables") or [] if str(item.get("text") or "").strip()).strip(),
        str(slide.get("notes") or "").strip(),
        str(slide.get("ocr_text") or "").strip(),
    ]
    deduped: List[str] = []
    seen = set()
    for part in ordered_parts:
        compact = " ".join(part.split()).strip()
        if not compact or compact in seen:
            continue
        seen.add(compact)
        deduped.append(part.strip())
    return "\n\n".join(deduped).strip()


def _build_slide_chunks(slides: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    chunks: List[Dict[str, Any]] = []
    for slide in slides:
        slide_text = _build_slide_text(slide)
        if not slide_text:
            continue
        slide_chunks = build_text_chunks(
            slide_text,
            kind="ppt_slide",
            extra={"slide_index": int(slide.get("slide_index") or 0), "title": str(slide.get("title") or "")},
        )
        for offset, chunk in enumerate(slide_chunks, start=1):
            chunk["index"] = len(chunks) + 1
            chunk["block_id"] = f"slide{slide.get('slide_index')}_chunk{offset}"
            chunks.append(chunk)
    return chunks


def _summarize_slides(slides: List[Dict[str, Any]]) -> str:
    preview_parts = []
    for slide in slides[:5]:
        title = str(slide.get("title") or "").strip() or f"第{slide.get('slide_index', '?')}页"
        body = _build_slide_text(slide)
        preview_parts.append(f"{title}: {clip_text(body, 80) if body else '无文本'}")
    if not preview_parts:
        return "PPT 中未提取到可用文本。"
    return "；".join(preview_parts)
