import os
import sys
import tempfile
import unittest


PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
if PROJECT_ROOT not in sys.path:
    sys.path.insert(0, PROJECT_ROOT)

from docx import Document  # noqa: E402
from pptx import Presentation  # noqa: E402

from backend.rag.docx_generation import render_docx, resolve_word_template  # noqa: E402
from backend.rag.ppt_generation import (  # noqa: E402
    build_theme_config,
    normalize_ppt_outline_to_12_slides,
    render_pptx,
    resolve_ppt_template,
)


class LessonPlanTemplateTests(unittest.TestCase):
    def test_official_word_template_renders_from_anchor_template(self):
        spec = {
            "requirement_summary": {
                "topic": "链表插入设计",
                "grade_subject": "数据结构",
                "duration": "45分钟",
                "teaching_goals": ["理解链表插入步骤", "掌握指针更新顺序"],
                "key_points": ["头插法", "尾插法"],
                "difficult_points": ["避免断链"],
                "knowledge_points": ["节点结构", "头插法", "尾插法"],
            },
            "docx_outline": [
                {
                    "section_title": "教学目标",
                    "section_goal": "目标聚焦",
                    "bullets": ["理解链表插入步骤"],
                    "source_refs": ["integration_source.txt"],
                }
            ],
        }
        core_spec = {
            "teaching_objectives": {
                "goals": ["理解链表插入步骤", "掌握指针更新顺序"],
                "key_points": ["头插法", "尾插法"],
                "difficult_points": ["避免断链"],
            },
            "student_profile": {
                "grade": "大一",
                "foundation": "已学指针基础",
                "learning_preference": "案例驱动",
                "common_misconceptions": ["容易遗漏后继节点保存"],
            },
            "knowledge_structure": {
                "knowledge_points": ["节点结构", "头插法", "尾插法"],
                "examples_or_cases": ["通讯录新增联系人"],
            },
            "teaching_flow": [
                {
                    "title": "案例导入",
                    "goal": "引入问题",
                    "activities": ["展示通讯录案例"],
                    "teacher_actions": ["抛出问题"],
                    "student_actions": ["回答"],
                    "assessment": ["口头提问"],
                },
                {
                    "title": "方法讲解",
                    "goal": "掌握步骤",
                    "activities": ["板书推演"],
                    "teacher_actions": ["讲解"],
                    "student_actions": ["记录"],
                    "assessment": ["判断题"],
                },
            ],
            "assessment_plan": {
                "questions": ["为什么要先保存后继节点？"],
                "in_class_checks": ["通过判断题定位错误"],
                "homework": ["完成伪代码整理"],
                "extension_tasks": ["比较顺序表与链表"],
            },
            "source_grounding": [{"claim": "先保存后继节点"}],
        }
        template_meta = resolve_word_template("official")
        self.assertTrue(template_meta["template_path"])

        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as temp_file:
            output_path = temp_file.name
        try:
            stats = render_docx(spec, output_path, core_spec=core_spec, template_profile="official")
            self.assertTrue(stats["template_used"])
            document = Document(output_path)
            joined = "\n".join([p.text for p in document.paragraphs if p.text.strip()])
            table_text = "\n".join(
                cell.text
                for table in document.tables
                for row in table.rows
                for cell in row.cells
                if cell.text.strip()
            )
            self.assertIn("链表插入设计", joined)
            self.assertIn("基础作业：完成伪代码整理", joined)
            self.assertIn("案例导入", table_text)
        finally:
            if os.path.exists(output_path):
                os.remove(output_path)

    def test_ppt_render_restores_blank_presentation_flow(self):
        spec = {
            "requirement_summary": {
                "topic": "链表插入设计",
                "grade_subject": "数据结构",
                "duration": "45分钟",
                "teaching_goals": ["理解链表插入步骤"],
                "key_points": ["头插法"],
                "knowledge_points": ["节点结构", "头插法"],
            },
            "ppt_outline": [
                {
                    "slide_type": "content",
                    "title": "案例导入",
                    "goal": "引入问题",
                    "bullets": ["展示通讯录案例", "对比顺序表和链表"],
                    "visual_suggestion": "案例流程图",
                    "source_refs": ["integration_source.txt"],
                },
                {
                    "slide_type": "summary",
                    "title": "课堂小结",
                    "goal": "",
                    "bullets": ["回顾关键步骤"],
                    "visual_suggestion": "",
                    "source_refs": ["integration_source.txt"],
                },
            ],
            "docx_outline": [],
        }
        template_meta = resolve_ppt_template("official")
        self.assertIsNone(template_meta["template_path"])
        self.assertEqual(template_meta["target_slide_count"], 12)
        plan = normalize_ppt_outline_to_12_slides(spec, template_meta.get("target_slide_count", 12))

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
            output_path = temp_file.name
        try:
            stats = render_pptx(
                plan,
                spec,
                build_theme_config(spec),
                [None] * len(plan["content_slides"]),
                output_path,
                template_profile="official",
            )
            self.assertNotIn("template_used", stats)
            presentation = Presentation(output_path)
            self.assertEqual(len(presentation.slides), 12)
            joined = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        joined.append(shape.text)
            text = "\n".join(joined)
            self.assertIn("链表插入设计", text)
            self.assertIn("课堂小结", text)
        finally:
            if os.path.exists(output_path):
                os.remove(output_path)


if __name__ == "__main__":
    unittest.main()
