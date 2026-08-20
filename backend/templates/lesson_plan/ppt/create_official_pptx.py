"""
Generate official.pptx -- a slide-master template with 7 custom slide layouts.

Layouts (0-based index among the CUSTOM layouts, after built-in ones):
  cover            -- ctrTitle + subTitle
  toc              -- title + body
  content_text     -- title + body (text-only)
  content_img_left -- title + picture(left) + body(right)
  content_img_right-- title + body(left) + picture(right)
  summary          -- title + body
  ending           -- ctrTitle + subTitle

Run:
    python create_official_pptx.py
"""

import copy
import os

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from pptx.oxml.ns import qn
from pptx.util import Inches
from pptx.dml.color import RGBColor

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT = "Microsoft YaHei"
TEXT_CLR = "252F40"
MUTED    = "627089"
ACCENT   = "2863D2"

OUT_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "official.pptx")


def _emu(inches_val):
    return int(inches_val)


def _make_ph_sp(ph_type, idx, name, left, top, width, height,
                font_size=None, bold=False, color=None, alignment=None):
    """Build a complete <p:sp> element for a placeholder."""
    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"

    sp = etree.Element(qn("p:sp"))

    # -- nvSpPr --
    nv = etree.SubElement(sp, qn("p:nvSpPr"))
    cNvPr = etree.SubElement(nv, qn("p:cNvPr"))
    cNvPr.set("id", str(idx + 10))
    cNvPr.set("name", name)
    cNvSpPr = etree.SubElement(nv, qn("p:cNvSpPr"))
    spLocks = etree.SubElement(cNvSpPr, qn("a:spLocks"))
    spLocks.set("noGrp", "1")
    nvPr = etree.SubElement(nv, qn("p:nvPr"))
    ph = etree.SubElement(nvPr, qn("p:ph"))
    if ph_type != "body":
        ph.set("type", ph_type)
    ph.set("idx", str(idx))

    # -- spPr --
    spPr = etree.SubElement(sp, qn("p:spPr"))
    xfrm = etree.SubElement(spPr, qn("a:xfrm"))
    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", str(_emu(left)))
    off.set("y", str(_emu(top)))
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", str(_emu(width)))
    ext.set("cy", str(_emu(height)))
    prstGeom = etree.SubElement(spPr, qn("a:prstGeom"))
    prstGeom.set("prst", "rect")
    etree.SubElement(prstGeom, qn("a:avLst"))

    # -- txBody --
    txBody = etree.SubElement(sp, qn("p:txBody"))
    bodyPr = etree.SubElement(txBody, qn("a:bodyPr"))
    bodyPr.set("wrap", "square")
    bodyPr.set("rtlCol", "0")
    etree.SubElement(txBody, qn("a:lstStyle"))
    p = etree.SubElement(txBody, qn("a:p"))

    if alignment:
        pPr = etree.SubElement(p, qn("a:pPr"))
        pPr.set("algn", alignment)

    r = etree.SubElement(p, qn("a:r"))
    rPr = etree.SubElement(r, qn("a:rPr"))
    rPr.set("lang", "zh-CN")
    rPr.set("altLang", "en-US")
    if bold:
        rPr.set("b", "1")
    if font_size:
        rPr.set("sz", str(int(font_size * 100)))
    if color:
        solidFill = etree.SubElement(rPr, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", color)
    latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", FONT)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", FONT)
    t = etree.SubElement(r, qn("a:t"))
    t.text = name

    return sp


def _make_pic_ph(idx, name, left, top, width, height):
    """Build a <p:sp> for a picture placeholder."""
    sp = etree.Element(qn("p:sp"))

    nv = etree.SubElement(sp, qn("p:nvSpPr"))
    cNvPr = etree.SubElement(nv, qn("p:cNvPr"))
    cNvPr.set("id", str(idx + 10))
    cNvPr.set("name", name)
    cNvSpPr = etree.SubElement(nv, qn("p:cNvSpPr"))
    spLocks = etree.SubElement(cNvSpPr, qn("a:spLocks"))
    spLocks.set("noGrp", "1")
    nvPr = etree.SubElement(nv, qn("p:nvPr"))
    ph = etree.SubElement(nvPr, qn("p:ph"))
    ph.set("type", "pic")
    ph.set("idx", str(idx))

    spPr = etree.SubElement(sp, qn("p:spPr"))
    xfrm = etree.SubElement(spPr, qn("a:xfrm"))
    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", str(_emu(left)))
    off.set("y", str(_emu(top)))
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", str(_emu(width)))
    ext.set("cy", str(_emu(height)))
    prstGeom = etree.SubElement(spPr, qn("a:prstGeom"))
    prstGeom.set("prst", "rect")
    etree.SubElement(prstGeom, qn("a:avLst"))

    txBody = etree.SubElement(sp, qn("p:txBody"))
    bodyPr = etree.SubElement(txBody, qn("a:bodyPr"))
    etree.SubElement(txBody, qn("a:lstStyle"))
    p = etree.SubElement(txBody, qn("a:p"))
    endParaRPr = etree.SubElement(p, qn("a:endParaRPr"))
    endParaRPr.set("lang", "zh-CN")

    return sp


def _build_layout_xml(layout_name, placeholders, base_xml):
    """Clone *base_xml* (blank layout) and replace its shapes with *placeholders*."""
    layout_xml = copy.deepcopy(base_xml)
    cSld = layout_xml.find(qn("p:cSld"))
    cSld.set("name", layout_name)

    spTree = cSld.find(qn("p:spTree"))
    for sp in list(spTree.findall(qn("p:sp"))):
        spTree.remove(sp)

    for ph in placeholders:
        if ph.get("is_pic"):
            sp = _make_pic_ph(
                idx=ph["idx"], name=ph["name"],
                left=ph["left"], top=ph["top"],
                width=ph["width"], height=ph["height"],
            )
        else:
            sp = _make_ph_sp(
                ph_type=ph["ph_type"], idx=ph["idx"], name=ph["name"],
                left=ph["left"], top=ph["top"],
                width=ph["width"], height=ph["height"],
                font_size=ph.get("font_size"),
                bold=ph.get("bold", False),
                color=ph.get("color"),
                alignment=ph.get("alignment"),
            )
        spTree.append(sp)

    return layout_xml


LAYOUT_SPECS = [
    {
        "name": "cover",
        "placeholders": [
            {"ph_type": "ctrTitle", "idx": 0, "name": "Title",
             "left": Inches(0.9), "top": Inches(1.7), "width": Inches(11.0), "height": Inches(1.4),
             "font_size": 30, "bold": True, "color": TEXT_CLR, "alignment": "l"},
            {"ph_type": "subTitle", "idx": 1, "name": "Subtitle",
             "left": Inches(0.95), "top": Inches(3.2), "width": Inches(8.8), "height": Inches(0.8),
             "font_size": 18, "color": MUTED, "alignment": "l"},
        ],
    },
    {
        "name": "toc",
        "placeholders": [
            {"ph_type": "title", "idx": 0, "name": "Title",
             "left": Inches(0.9), "top": Inches(0.55), "width": Inches(8.9), "height": Inches(0.7),
             "font_size": 28, "bold": True, "color": TEXT_CLR},
            {"ph_type": "body", "idx": 1, "name": "TOC Body",
             "left": Inches(0.9), "top": Inches(1.6), "width": Inches(11.4), "height": Inches(5.0),
             "font_size": 18, "color": TEXT_CLR},
        ],
    },
    {
        "name": "content_text",
        "placeholders": [
            {"ph_type": "title", "idx": 0, "name": "Title",
             "left": Inches(0.9), "top": Inches(0.55), "width": Inches(8.9), "height": Inches(0.7),
             "font_size": 28, "bold": True, "color": TEXT_CLR},
            {"ph_type": "body", "idx": 1, "name": "Content Body",
             "left": Inches(0.95), "top": Inches(1.55), "width": Inches(11.35), "height": Inches(4.5),
             "font_size": 17, "color": TEXT_CLR},
        ],
    },
    {
        "name": "content_img_left",
        "placeholders": [
            {"ph_type": "title", "idx": 0, "name": "Title",
             "left": Inches(0.9), "top": Inches(0.55), "width": Inches(8.9), "height": Inches(0.7),
             "font_size": 28, "bold": True, "color": TEXT_CLR},
            {"ph_type": "pic", "idx": 2, "name": "Picture Left",
             "left": Inches(0.7), "top": Inches(1.55), "width": Inches(5.0), "height": Inches(4.2),
             "is_pic": True},
            {"ph_type": "body", "idx": 1, "name": "Content Right",
             "left": Inches(6.0), "top": Inches(1.55), "width": Inches(6.35), "height": Inches(4.5),
             "font_size": 17, "color": TEXT_CLR},
        ],
    },
    {
        "name": "content_img_right",
        "placeholders": [
            {"ph_type": "title", "idx": 0, "name": "Title",
             "left": Inches(0.9), "top": Inches(0.55), "width": Inches(8.9), "height": Inches(0.7),
             "font_size": 28, "bold": True, "color": TEXT_CLR},
            {"ph_type": "body", "idx": 1, "name": "Content Left",
             "left": Inches(0.8), "top": Inches(1.55), "width": Inches(6.2), "height": Inches(4.5),
             "font_size": 17, "color": TEXT_CLR},
            {"ph_type": "pic", "idx": 2, "name": "Picture Right",
             "left": Inches(7.6), "top": Inches(1.55), "width": Inches(5.0), "height": Inches(4.2),
             "is_pic": True},
        ],
    },
    {
        "name": "summary",
        "placeholders": [
            {"ph_type": "title", "idx": 0, "name": "Title",
             "left": Inches(0.9), "top": Inches(0.55), "width": Inches(8.9), "height": Inches(0.7),
             "font_size": 28, "bold": True, "color": TEXT_CLR},
            {"ph_type": "body", "idx": 1, "name": "Summary Body",
             "left": Inches(0.9), "top": Inches(1.6), "width": Inches(11.4), "height": Inches(5.0),
             "font_size": 17, "color": TEXT_CLR},
        ],
    },
    {
        "name": "ending",
        "placeholders": [
            {"ph_type": "ctrTitle", "idx": 0, "name": "Title",
             "left": Inches(1.2), "top": Inches(2.2), "width": Inches(10.8), "height": Inches(1.1),
             "font_size": 30, "bold": True, "color": TEXT_CLR, "alignment": "ctr"},
            {"ph_type": "subTitle", "idx": 1, "name": "Subtitle",
             "left": Inches(1.6), "top": Inches(3.45), "width": Inches(10.0), "height": Inches(0.7),
             "font_size": 18, "color": MUTED, "alignment": "ctr"},
        ],
    },
]


def build_template():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    master = prs.slide_masters[0]
    master_part = master.part
    master_elem = master_part._element

    blank_layout = prs.slide_layouts[6]
    blank_layout_xml = copy.deepcopy(blank_layout.part._element)

    sldLayoutIdLst = master_elem.find(qn("p:sldLayoutIdLst"))
    max_id = max(
        (int(el.get("id", "0")) for el in sldLayoutIdLst.findall(qn("p:sldLayoutId"))),
        default=2147483648,
    )

    existing_rels = [
        rel.target_partname
        for rel in master_part.rels.values()
        if "slideLayout" in str(rel.target_partname)
    ]
    next_layout_num = len(existing_rels) + 1

    added = []

    for spec in LAYOUT_SPECS:
        max_id += 1
        layout_xml = _build_layout_xml(spec["name"], spec["placeholders"], blank_layout_xml)

        partname = PackURI(f"/ppt/slideLayouts/slideLayout{next_layout_num}.xml")
        content_type = (
            "application/vnd.openxmlformats-officedocument"
            ".presentationml.slideLayout+xml"
        )
        xml_bytes = etree.tostring(layout_xml, xml_declaration=True, encoding="UTF-8", standalone=True)
        new_part = Part(partname, content_type, master_part.package, xml_bytes)

        rId = master_part.relate_to(new_part, RT.SLIDE_LAYOUT)
        new_part.relate_to(master_part, RT.SLIDE_MASTER)

        sldLayoutId = etree.SubElement(sldLayoutIdLst, qn("p:sldLayoutId"))
        sldLayoutId.set("id", str(max_id))
        sldLayoutId.set(qn("r:id"), rId)

        added.append(f"  [{next_layout_num - 1}] {spec['name']}")
        next_layout_num += 1

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")
    print(f"Added {len(added)} custom layouts:")
    for line in added:
        print(line)

    verify_prs = Presentation(OUT_PATH)
    print(f"\nVerification ({len(verify_prs.slide_layouts)} total layouts):")
    for i, lay in enumerate(verify_prs.slide_layouts):
        print(f"  [{i}] {lay.name}")


if __name__ == "__main__":
    build_template()
